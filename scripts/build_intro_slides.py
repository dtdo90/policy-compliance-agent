"""Generate the 3 intro slides for the Travel Agency Policy Studio demo.

Run:
    uv run python scripts/build_intro_slides.py

Output: reference/presentation/intro_slides.pptx

The styling intentionally mirrors the plain reference deck: white background,
clean Calibri text, navy headings, a single blue accent for arrows/dividers.
No hero bands, no gradients, no colored pills.
"""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

NAVY = RGBColor(0x1F, 0x38, 0x64)
INK = RGBColor(0x33, 0x33, 0x33)
GREY = RGBColor(0x59, 0x59, 0x59)
ACCENT = RGBColor(0x2E, 0x75, 0xB6)
DIVIDER = RGBColor(0xBF, 0xBF, 0xBF)
SOFT = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=16,
    bold=False,
    italic=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.2,
    font="Calibri",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def add_runs(slide, left, top, width, height, paragraphs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.25):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    for i, paragraph in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for spec in paragraph:
            run = p.add_run()
            run.text = spec["text"]
            run.font.name = spec.get("font", "Calibri")
            run.font.size = Pt(spec.get("size", 16))
            run.font.bold = spec.get("bold", False)
            run.font.italic = spec.get("italic", False)
            run.font.color.rgb = spec.get("color", INK)
    return box


def add_divider(slide, left, top, width, *, color=DIVIDER, weight=0.75):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_box(slide, left, top, width, height, *, fill=WHITE, border=NAVY, border_w=1.0, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if rounded:
        shape.adjustments[0] = 0.10
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border
        shape.line.width = Pt(border_w)
    shape.shadow.inherit = False
    return shape


def add_line(slide, x1, y1, x2, y2, *, color=ACCENT, weight=1.5, head_at_end=False):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    if head_at_end:
        ln = line.line._get_or_add_ln()
        for existing in ln.findall(qn("a:tailEnd")):
            ln.remove(existing)
        tail = etree.SubElement(ln, qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", "med")
        tail.set("len", "med")
    return line


def add_arrow(slide, left, top, width, height, *, color=ACCENT, direction="right"):
    shape_type = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
    }[direction]
    arrow = slide.shapes.add_shape(shape_type, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow


def add_slide_title(slide, eyebrow, title):
    add_text(slide, Inches(0.5), Inches(0.30), Inches(12.3), Inches(0.30), eyebrow, size=12, bold=True, color=ACCENT)
    add_text(slide, Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.55), title, size=28, bold=True, color=NAVY)
    add_divider(slide, Inches(0.5), Inches(1.10), Inches(12.33), color=NAVY, weight=1.25)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ---------------------------------------------------------------------------
# Slide 1 - Title + the two rules
# ---------------------------------------------------------------------------
s1 = prs.slides.add_slide(blank)
add_text(
    s1,
    Inches(0.5),
    Inches(0.65),
    Inches(12.3),
    Inches(0.40),
    "UOB AI LABS  |  PBVA",
    size=12,
    bold=True,
    color=ACCENT,
)
add_text(
    s1,
    Inches(0.5),
    Inches(0.95),
    Inches(12.3),
    Inches(0.95),
    "Travel Agency Policy Studio",
    size=40,
    bold=True,
    color=NAVY,
)
add_text(
    s1,
    Inches(0.5),
    Inches(1.85),
    Inches(12.3),
    Inches(0.45),
    "AI compliance for travel-agency conversations - two rules, two anchors.",
    size=18,
    color=GREY,
)
add_divider(s1, Inches(0.5), Inches(2.40), Inches(12.33), color=DIVIDER, weight=1.0)

# Rule 101 - Account Unlock
r1_top = Inches(2.65)
add_text(s1, Inches(0.5), r1_top, Inches(12.3), Inches(0.45), "Rule 101 - Account Unlock", size=22, bold=True, color=NAVY)
add_text(
    s1,
    Inches(0.5),
    r1_top + Inches(0.50),
    Inches(12.3),
    Inches(0.40),
    "Verify the customer's identity before resetting or unlocking the account.",
    size=15,
    color=INK,
)
add_runs(
    s1,
    Inches(0.5),
    r1_top + Inches(0.92),
    Inches(12.3),
    Inches(0.45),
    [
        [
            {"text": "Anchor: ", "size": 14, "bold": True, "color": ACCENT},
            {"text": "\u201cBefore I reset or unlock your account, I need to verify your identity first.\u201d", "size": 14, "color": INK, "italic": True},
        ]
    ],
)

# Rule 102 - Flight Change
r2_top = Inches(4.40)
add_text(s1, Inches(0.5), r2_top, Inches(12.3), Inches(0.45), "Rule 102 - Flight Change", size=22, bold=True, color=NAVY)
add_text(
    s1,
    Inches(0.5),
    r2_top + Inches(0.50),
    Inches(12.3),
    Inches(0.40),
    "Disclose the change fee AND the fare difference / travel credit before confirming the change.",
    size=15,
    color=INK,
)
add_runs(
    s1,
    Inches(0.5),
    r2_top + Inches(0.92),
    Inches(12.3),
    Inches(0.45),
    [
        [
            {"text": "Anchor 1: ", "size": 14, "bold": True, "color": ACCENT},
            {"text": "\u201cBefore I confirm this booking change, there is a change fee that will apply.\u201d", "size": 14, "color": INK, "italic": True},
        ]
    ],
)
add_runs(
    s1,
    Inches(0.5),
    r2_top + Inches(1.30),
    Inches(12.3),
    Inches(0.55),
    [
        [
            {"text": "Anchor 2: ", "size": 14, "bold": True, "color": ACCENT},
            {"text": "\u201cThere is also a fare difference on the new itinerary - you will either pay the extra or receive the balance as travel credit.\u201d", "size": 14, "color": INK, "italic": True},
        ]
    ],
)


# ---------------------------------------------------------------------------
# Slide 2 - Inference pipeline
# ---------------------------------------------------------------------------
s2 = prs.slides.add_slide(blank)
add_slide_title(s2, "PART 1", "Inference pipeline")

steps = [
    ("1. Chunk", "Split transcript into\nshort utterances."),
    ("2. Encode", "Sentence Transformer\nmaps each chunk\nto a vector."),
    ("3. Retrieve", "Top-k chunks closest\nto the rule's anchor."),
    ("4. Rerank", "Cross Encoder scores\neach (anchor, chunk)\npair from 0 to 1."),
    ("5. Decide", "Take MAX score.\n  > 0.5 -> compliant\n  <= 0.5 -> missing"),
]

n = len(steps)
step_w = Inches(2.20)
step_h = Inches(1.85)
gap = Inches(0.20)
total_w = n * step_w + (n - 1) * gap
start_x = (prs.slide_width - total_w) // 2
step_y = Inches(2.10)

for i, (head, body) in enumerate(steps):
    left = start_x + i * (step_w + gap)
    add_box(s2, left, step_y, step_w, step_h, fill=WHITE, border=NAVY, border_w=1.25)
    add_text(s2, left, step_y + Inches(0.15), step_w, Inches(0.45), head, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s2, left + Inches(0.10), step_y + Inches(0.62), step_w - Inches(0.20), step_h - Inches(0.70), body, size=12, color=INK, align=PP_ALIGN.CENTER)
    if i < n - 1:
        ax = left + step_w + Inches(0.02)
        ay = step_y + step_h / 2 - Inches(0.13)
        add_arrow(s2, ax, ay, gap - Inches(0.04), Inches(0.26), color=ACCENT)

# Models row
mod_top = step_y + step_h + Inches(0.45)
add_runs(
    s2,
    Inches(0.5),
    mod_top,
    Inches(12.3),
    Inches(0.40),
    [
        [
            {"text": "Models:  ", "size": 14, "bold": True, "color": NAVY},
            {"text": "Sentence Transformer (MiniLM-L6, retriever)  +  Cross Encoder (MiniLM-L12, verifier)", "size": 14, "color": INK},
        ]
    ],
    align=PP_ALIGN.CENTER,
)

# Score band
band_top = mod_top + Inches(0.70)
band_left = Inches(1.0)
band_w = Inches(11.33)
zone_h = Inches(0.34)

fail_w = Inches(11.33 * 0.30)
border_w = Inches(11.33 * 0.40)
pass_w = Inches(11.33 * 0.30)

fail_zone = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, band_left, band_top, fail_w, zone_h)
fail_zone.fill.solid(); fail_zone.fill.fore_color.rgb = SOFT
fail_zone.line.color.rgb = DIVIDER; fail_zone.line.width = Pt(0.5)
border_zone = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, band_left + fail_w, band_top, border_w, zone_h)
border_zone.fill.solid(); border_zone.fill.fore_color.rgb = RGBColor(0xFF, 0xF1, 0xCE)
border_zone.line.color.rgb = DIVIDER; border_zone.line.width = Pt(0.5)
pass_zone = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, band_left + fail_w + border_w, band_top, pass_w, zone_h)
pass_zone.fill.solid(); pass_zone.fill.fore_color.rgb = SOFT
pass_zone.line.color.rgb = DIVIDER; pass_zone.line.width = Pt(0.5)

add_text(s2, band_left, band_top + Inches(0.40), fail_w, Inches(0.30), "fail (\u2264 0.30)", size=11, color=GREY, align=PP_ALIGN.CENTER)
add_text(s2, band_left + fail_w, band_top + Inches(0.40), border_w, Inches(0.30), "borderline (0.30 - 0.70)  -  human review", size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s2, band_left + fail_w + border_w, band_top + Inches(0.40), pass_w, Inches(0.30), "pass (> 0.50)", size=11, color=GREY, align=PP_ALIGN.CENTER)

add_text(s2, band_left - Inches(0.15), band_top - Inches(0.30), Inches(0.5), Inches(0.25), "0.0", size=10, color=GREY, align=PP_ALIGN.CENTER)
add_text(s2, band_left + band_w - Inches(0.35), band_top - Inches(0.30), Inches(0.5), Inches(0.25), "1.0", size=10, color=GREY, align=PP_ALIGN.CENTER)

# Caption
add_text(
    s2,
    Inches(0.5),
    band_top + Inches(0.95),
    Inches(12.3),
    Inches(0.40),
    "The same pipeline runs once per rule, with that rule's anchor.",
    size=14,
    italic=True,
    color=GREY,
    align=PP_ALIGN.CENTER,
)


# ---------------------------------------------------------------------------
# Slide 3 - Agentic workflow (runs locally)
# ---------------------------------------------------------------------------
s3 = prs.slides.add_slide(blank)
add_slide_title(s3, "PART 2", "Agentic workflow  -  runs locally")

# Two-column layout: HITL (left) -> Agentic (right). Columns sit a little
# inboard from the slide edges so the loop arrows + supervisor label stay
# on-slide.
col_top = Inches(1.40)
col_w = Inches(5.05)
gap = Inches(0.20)
left_col = Inches(0.95)
right_col = left_col + col_w + gap

# Column titles
add_text(s3, left_col, col_top, col_w, Inches(0.40), "Human-in-the-loop  (today)", size=15, bold=True, color=GREY, align=PP_ALIGN.CENTER)
add_text(s3, right_col, col_top, col_w, Inches(0.40), "Agentic  (this demo)", size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# HITL nodes
hitl_nodes = [
    "Synthetic data",
    "Train models\n(retriever + verifier)",
    "Inference\n(input transcripts)",
    "Human review\n(pass + borderline)",
    "Data augmentation\n(fix gaps)",
    "Retrain models",
]

# Agentic nodes
agentic_nodes = [
    ("Synthetic data", None),
    ("TrainerAgent", "retriever + verifier"),
    ("InferenceAgent", "scores per (rule, transcript)"),
    ("ReviewAgent + Human", "approve borderline / fail"),
    ("InvestigatorAgent", "diagnose gaps, generate variants"),
    ("TrainerAgent", "retrain + evaluate new metrics"),
]

content_top = col_top + Inches(0.50)
node_h = Inches(0.50)
node_gap = Inches(0.08)
node_w = col_w - Inches(0.50)
inference_idx = 2  # 0=Synthetic, 1=Train, 2=Inference, ...
last_idx = 5  # last node (Retrain / TrainerAgent retrain)

for i, label in enumerate(hitl_nodes):
    y = content_top + i * (node_h + node_gap)
    add_box(s3, left_col + Inches(0.25), y, node_w, node_h, fill=WHITE, border=GREY, border_w=0.75)
    lines = label.split("\n")
    if len(lines) == 1:
        add_text(s3, left_col + Inches(0.25), y, node_w, node_h, lines[0], size=12, bold=True, color=GREY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    else:
        add_text(s3, left_col + Inches(0.25), y + Inches(0.05), node_w, Inches(0.28), lines[0], size=12, bold=True, color=GREY, align=PP_ALIGN.CENTER)
        add_text(s3, left_col + Inches(0.25), y + Inches(0.30), node_w, Inches(0.25), lines[1], size=10, color=GREY, align=PP_ALIGN.CENTER)
    if i < len(hitl_nodes) - 1:
        ay = y + node_h + Inches(0.005)
        ax = left_col + col_w / 2 - Inches(0.08)
        add_arrow(s3, ax, ay, Inches(0.16), node_gap - Inches(0.01), color=GREY, direction="down")

# Right column - agentic
for i, (title, sub) in enumerate(agentic_nodes):
    y = content_top + i * (node_h + node_gap)
    border = NAVY
    add_box(s3, right_col + Inches(0.25), y, node_w, node_h, fill=WHITE, border=border, border_w=1.0)
    if sub is None:
        add_text(s3, right_col + Inches(0.25), y, node_w, node_h, title, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    else:
        add_text(s3, right_col + Inches(0.25), y + Inches(0.04), node_w, Inches(0.26), title, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s3, right_col + Inches(0.25), y + Inches(0.28), node_w, Inches(0.25), sub, size=10, color=INK, align=PP_ALIGN.CENTER)
    if i < len(agentic_nodes) - 1:
        ay = y + node_h + Inches(0.005)
        ax = right_col + col_w / 2 - Inches(0.08)
        add_arrow(s3, ax, ay, Inches(0.16), node_gap - Inches(0.01), color=ACCENT, direction="down")

# Loop arrow: bottom of last node -> back up to right side of "Inference" node.
# Drawn as three straight connectors with an arrowhead on the final segment.

def _draw_loop(slide, *, col_left, col_w_, side, color):
    """Draw a U-shaped loop arrow from the last node back up to inference."""
    last_y_mid = content_top + last_idx * (node_h + node_gap) + node_h / 2
    inf_y_mid = content_top + inference_idx * (node_h + node_gap) + node_h / 2
    box_left = col_left + Inches(0.25)
    box_right = box_left + node_w
    if side == "right":
        track_x = box_right + Inches(0.22)
        anchor_x = box_right
    else:
        track_x = box_left - Inches(0.22)
        anchor_x = box_left
    add_line(slide, anchor_x, last_y_mid, track_x, last_y_mid, color=color, weight=1.5)
    add_line(slide, track_x, last_y_mid, track_x, inf_y_mid, color=color, weight=1.5)
    add_line(slide, track_x, inf_y_mid, anchor_x, inf_y_mid, color=color, weight=1.5, head_at_end=True)
    return track_x, last_y_mid, inf_y_mid


# HITL loop on the left side of the left column (grey)
_draw_loop(s3, col_left=left_col, col_w_=col_w, side="left", color=GREY)

# Agentic loop on the right side of the right column (accent blue)
track_x_r, last_y_mid_r, inf_y_mid_r = _draw_loop(
    s3, col_left=right_col, col_w_=col_w, side="right", color=ACCENT
)

# SupervisorAgent label, placed to the right of the agentic loop, vertically
# centered between the loop's top and bottom horizontal segments.
sup_mid_y = (last_y_mid_r + inf_y_mid_r) / 2
add_text(
    s3,
    track_x_r + Inches(0.12),
    sup_mid_y - Inches(0.30),
    Inches(1.55),
    Inches(0.32),
    "SupervisorAgent",
    size=12,
    bold=True,
    color=NAVY,
)
add_text(
    s3,
    track_x_r + Inches(0.12),
    sup_mid_y + Inches(0.02),
    Inches(1.55),
    Inches(0.50),
    "orchestrates the\nfull loop",
    size=10,
    color=GREY,
    line_spacing=1.15,
)

# Bottom callout: runs locally
local_top = Inches(6.85)
add_box(s3, Inches(0.5), local_top, Inches(12.33), Inches(0.45), fill=SOFT, border=NAVY, border_w=0.75)
add_runs(
    s3,
    Inches(0.7),
    local_top + Inches(0.07),
    Inches(12.0),
    Inches(0.32),
    [
        [
            {"text": "Runs locally:  ", "size": 13, "bold": True, "color": NAVY},
            {"text": "MiniLM models on CPU + a local LLM (Ollama) on the same laptop.  No external API required.", "size": 13, "color": INK},
        ]
    ],
)


# ---------------------------------------------------------------------------
# Speaker notes
# ---------------------------------------------------------------------------
NOTES = [
    (
        s1,
        "Slide 1 - Travel Agency Policy Studio (~1 min).\n\n"
        "Today's demo: AI that checks whether a travel-agency agent said the right thing on the call.\n"
        "Two rules in scope: (101) Account Unlock - verify identity before resetting or unlocking; "
        "(102) Flight Change - disclose change fee AND fare difference / travel credit before confirming.\n"
        "Each rule is grounded by a representative anchor sentence. Adding a new rule = adding a new anchor."
    ),
    (
        s2,
        "Slide 2 - Inference pipeline (~1.5 min).\n\n"
        "Five stages: chunk the transcript, encode chunks with the Sentence Transformer, retrieve top-k closest to the "
        "rule's anchor, rerank pairs with the Cross Encoder, then take the MAX score. Score > 0.5 = compliant, otherwise missing.\n"
        "Borderline 0.30 - 0.70 is sent to a human. The same pipeline runs once per rule."
    ),
    (
        s3,
        "Slide 3 - Agentic workflow (~1.5 min).\n\n"
        "Left: the original human-in-the-loop loop. Right: every manual step becomes a named agent - TrainerAgent, "
        "InferenceAgent, ReviewAgent + Human, InvestigatorAgent, TrainerAgent. SupervisorAgent "
        "orchestrates the full loop.\n"
        "Emphasis: this entire stack runs locally on a laptop - MiniLM models on CPU and a local LLM via Ollama. No external API."
    ),
]
for slide, text in NOTES:
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    notes_tf.text = text


out = Path("reference/presentation/intro_slides.pptx")
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f"Wrote {out}")
