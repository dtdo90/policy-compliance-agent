"""Replace the agentic-workflow slide in intro_slides.pptx with two slides:

  - Slide 3:  Human-in-the-loop today  (+ iteration objective)
  - Slide 4:  Agentic workflow (this demo, runs locally)

Slides 1 (title) and 2 (inference pipeline -- hand-edited by user) are left
untouched.

Run:
    uv run python scripts/update_agentic_slides.py
"""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

NAVY = RGBColor(0x1F, 0x38, 0x64)
INK = RGBColor(0x33, 0x33, 0x33)
GREY = RGBColor(0x59, 0x59, 0x59)
ACCENT = RGBColor(0x2E, 0x75, 0xB6)
DIVIDER = RGBColor(0xBF, 0xBF, 0xBF)
SOFT = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xC2, 0x83, 0x00)


# ---------------------------------------------------------------------------
# Helpers (duplicated from build_intro_slides.py so this script stands alone)
# ---------------------------------------------------------------------------
def add_text(slide, left, top, width, height, text, *, size=16, bold=False, italic=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def add_runs(slide, left, top, width, height, paragraphs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.25):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    for i, paragraph in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for spec in paragraph:
            run = p.add_run()
            run.text = spec["text"]
            run.font.name = spec.get("font", "Calibri")
            run.font.size = Pt(spec.get("size", 16))
            run.font.bold = spec.get("bold", False)
            run.font.italic = spec.get("italic", False)
            run.font.color.rgb = spec.get("color", INK)
    return box


def add_divider(slide, left, top, width, *, color=DIVIDER, weight=0.75):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_box(slide, left, top, width, height, *, fill=WHITE, border=NAVY, border_w=1.0, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if rounded:
        shape.adjustments[0] = 0.10
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border
        shape.line.width = Pt(border_w)
    shape.shadow.inherit = False
    return shape


def add_line(slide, x1, y1, x2, y2, *, color=ACCENT, weight=1.5, head_at_end=False):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    if head_at_end:
        ln = line.line._get_or_add_ln()
        for existing in ln.findall(qn("a:tailEnd")):
            ln.remove(existing)
        tail = etree.SubElement(ln, qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", "med"); tail.set("len", "med")
    return line


def add_arrow(slide, left, top, width, height, *, color=ACCENT, direction="right"):
    shape_type = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
    }[direction]
    arrow = slide.shapes.add_shape(shape_type, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow


def add_slide_title(slide, eyebrow, title):
    add_text(slide, Inches(0.5), Inches(0.30), Inches(12.3), Inches(0.30), eyebrow, size=12, bold=True, color=ACCENT)
    add_text(slide, Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.55), title, size=28, bold=True, color=NAVY)
    add_divider(slide, Inches(0.5), Inches(1.10), Inches(12.33), color=NAVY, weight=1.25)


def delete_slide(prs, index):
    """Remove a slide by index from a presentation, fixing up relationships."""
    xml_slides = prs.slides._sldIdLst
    slide_ids = list(xml_slides)
    slide_id = slide_ids[index]
    rId = slide_id.attrib[
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    ]
    prs.part.drop_rel(rId)
    xml_slides.remove(slide_id)


# ---------------------------------------------------------------------------
# Slide 3 - Human-in-the-loop today + iteration objective
# ---------------------------------------------------------------------------
def build_hitl_slide(prs):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_slide_title(s, "PART 2", "Human-in-the-loop  -  today")

    # Left: 6-step HITL flow with a U-shaped loop arrow on its left side
    col_top = Inches(1.40)
    col_w = Inches(5.40)
    left_col = Inches(0.95)

    add_text(
        s, left_col, col_top, col_w, Inches(0.40),
        "Workflow", size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
    )

    nodes = [
        ("Synthetic data", None),
        ("Train models", "retriever + verifier"),
        ("Inference", "input transcripts -> scores"),
        ("Human review", "approve pass + borderline"),
        ("Data augmentation", "fix gaps"),
        ("Retrain models", "observe movement"),
    ]
    inference_idx = 2
    last_idx = len(nodes) - 1

    content_top = col_top + Inches(0.50)
    node_h = Inches(0.55)
    node_gap = Inches(0.10)
    node_w = col_w - Inches(0.50)

    for i, (head, sub) in enumerate(nodes):
        y = content_top + i * (node_h + node_gap)
        add_box(s, left_col + Inches(0.25), y, node_w, node_h, fill=WHITE, border=GREY, border_w=0.85)
        if sub is None:
            add_text(s, left_col + Inches(0.25), y, node_w, node_h, head, size=13, bold=True, color=GREY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        else:
            add_text(s, left_col + Inches(0.25), y + Inches(0.05), node_w, Inches(0.28), head, size=13, bold=True, color=GREY, align=PP_ALIGN.CENTER)
            add_text(s, left_col + Inches(0.25), y + Inches(0.30), node_w, Inches(0.25), sub, size=10, color=GREY, align=PP_ALIGN.CENTER)
        if i < len(nodes) - 1:
            ay = y + node_h + Inches(0.005)
            ax = left_col + col_w / 2 - Inches(0.08)
            add_arrow(s, ax, ay, Inches(0.16), node_gap - Inches(0.01), color=GREY, direction="down")

    # U-loop on the LEFT side: from bottom of "Retrain" back up to right of "Inference"
    inf_mid_y = content_top + inference_idx * (node_h + node_gap) + node_h / 2
    last_mid_y = content_top + last_idx * (node_h + node_gap) + node_h / 2
    box_left = left_col + Inches(0.25)
    track_x = box_left - Inches(0.30)
    add_line(s, box_left, last_mid_y, track_x, last_mid_y, color=GREY, weight=1.5)
    add_line(s, track_x, last_mid_y, track_x, inf_mid_y, color=GREY, weight=1.5)
    add_line(s, track_x, inf_mid_y, box_left, inf_mid_y, color=GREY, weight=1.5, head_at_end=True)

    # Right panel: Iteration objective
    panel_left = Inches(7.10)
    panel_w = Inches(5.70)
    panel_top = Inches(1.50)
    panel_h = Inches(5.20)

    add_box(s, panel_left, panel_top, panel_w, panel_h, fill=SOFT, border=NAVY, border_w=1.0)
    add_text(
        s, panel_left + Inches(0.30), panel_top + Inches(0.18), panel_w - Inches(0.6), Inches(0.30),
        "ITERATION OBJECTIVE", size=12, bold=True, color=ACCENT,
    )
    add_text(
        s, panel_left + Inches(0.30), panel_top + Inches(0.46), panel_w - Inches(0.6), Inches(0.55),
        "What we want each loop pass to achieve",
        size=18, bold=True, color=NAVY,
    )
    add_divider(
        s, panel_left + Inches(0.30), panel_top + Inches(1.05),
        panel_w - Inches(0.6), color=DIVIDER, weight=0.75,
    )

    # Bullet 1
    b1_top = panel_top + Inches(1.25)
    bullet1 = s.shapes.add_shape(MSO_SHAPE.OVAL, panel_left + Inches(0.30), b1_top + Inches(0.10), Inches(0.18), Inches(0.18))
    bullet1.fill.solid(); bullet1.fill.fore_color.rgb = ACCENT; bullet1.line.fill.background()
    add_runs(
        s, panel_left + Inches(0.62), b1_top, panel_w - Inches(0.95), Inches(0.85),
        [
            [
                {"text": "Wrong predictions ", "size": 16, "color": INK},
                {"text": "become correct.", "size": 16, "color": NAVY, "bold": True},
            ]
        ],
    )

    # Bullet 2
    b2_top = b1_top + Inches(1.10)
    bullet2 = s.shapes.add_shape(MSO_SHAPE.OVAL, panel_left + Inches(0.30), b2_top + Inches(0.10), Inches(0.18), Inches(0.18))
    bullet2.fill.solid(); bullet2.fill.fore_color.rgb = ACCENT; bullet2.line.fill.background()
    add_runs(
        s, panel_left + Inches(0.62), b2_top, panel_w - Inches(0.95), Inches(1.20),
        [
            [
                {"text": "Borderline scores ", "size": 16, "color": INK},
                {"text": "(0.30 - 0.70)", "size": 16, "color": AMBER, "bold": True},
                {"text": " move closer to confident ", "size": 16, "color": INK},
                {"text": "0 or 1", "size": 16, "color": NAVY, "bold": True},
                {"text": ".", "size": 16, "color": INK},
            ]
        ],
    )

    # Mini visual under bullet 2: a score band with arrows pulling outward
    band_top = b2_top + Inches(1.55)
    band_w_in = 5.70 - 1.10  # = 4.60 inches
    band_left = panel_left + Inches(0.55)
    band_w = Inches(band_w_in)
    zone_h = Inches(0.30)
    fail_w = Inches(band_w_in * 0.30)
    border_w = Inches(band_w_in * 0.40)
    pass_w = Inches(band_w_in * 0.30)

    fail_zone = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, band_left, band_top, fail_w, zone_h)
    fail_zone.fill.solid(); fail_zone.fill.fore_color.rgb = WHITE
    fail_zone.line.color.rgb = DIVIDER; fail_zone.line.width = Pt(0.5)
    border_zone = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, band_left + fail_w, band_top, border_w, zone_h)
    border_zone.fill.solid(); border_zone.fill.fore_color.rgb = RGBColor(0xFF, 0xF1, 0xCE)
    border_zone.line.color.rgb = DIVIDER; border_zone.line.width = Pt(0.5)
    pass_zone = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, band_left + fail_w + border_w, band_top, pass_w, zone_h)
    pass_zone.fill.solid(); pass_zone.fill.fore_color.rgb = WHITE
    pass_zone.line.color.rgb = DIVIDER; pass_zone.line.width = Pt(0.5)

    add_text(s, band_left - Inches(0.20), band_top + Inches(0.32), Inches(0.5), Inches(0.25), "0", size=10, color=GREY, align=PP_ALIGN.CENTER)
    add_text(s, band_left + band_w - Inches(0.30), band_top + Inches(0.32), Inches(0.5), Inches(0.25), "1", size=10, color=GREY, align=PP_ALIGN.CENTER)

    # Two arrows, one from middle pulling left, one pulling right
    arrow_y = band_top + Inches(0.04)
    mid_x = band_left + fail_w + border_w / 2
    add_arrow(s, mid_x - Inches(0.50), arrow_y, Inches(0.40), Inches(0.22), color=ACCENT, direction="left")
    add_arrow(s, mid_x + Inches(0.10), arrow_y, Inches(0.40), Inches(0.22), color=ACCENT, direction="right")

    add_text(
        s, band_left, band_top + Inches(0.42), band_w, Inches(0.30),
        "borderline -> 0 or 1",
        size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER,
    )

    return s


# ---------------------------------------------------------------------------
# Slide 4 - Agentic workflow (this demo, runs locally)
# ---------------------------------------------------------------------------
def build_agentic_slide(prs):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_slide_title(s, "PART 3", "Agentic workflow  -  this demo, runs locally")

    # Centered single column of 6 agent boxes
    col_w = Inches(4.20)
    col_x = (prs.slide_width - col_w) / 2
    col_top = Inches(1.30)

    add_text(
        s, col_x, col_top, col_w, Inches(0.35),
        "Agents in the loop", size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
    )

    nodes = [
        ("Synthetic data", None),
        ("TrainerAgent", "retriever + verifier"),
        ("InferenceAgent", "scores per (rule, transcript)"),
        ("ReviewAgent + Human", "approve borderline / fail"),
        ("InvestigatorAgent", "diagnose gaps, generate variants"),
        ("TrainerAgent", "retrain + evaluate new metrics"),
    ]
    inference_idx = 2
    last_idx = len(nodes) - 1

    content_top = col_top + Inches(0.45)
    node_h = Inches(0.62)
    node_gap = Inches(0.12)
    node_w = col_w - Inches(0.40)

    for i, (head, sub) in enumerate(nodes):
        y = content_top + i * (node_h + node_gap)
        # Highlight the human-approval node faintly
        if i == 3:
            fill = RGBColor(0xFF, 0xF6, 0xE0)
            border = RGBColor(0xC2, 0x83, 0x00)
        else:
            fill = WHITE
            border = NAVY
        add_box(s, col_x + Inches(0.20), y, node_w, node_h, fill=fill, border=border, border_w=1.0)
        if sub is None:
            add_text(s, col_x + Inches(0.20), y, node_w, node_h, head, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        else:
            add_text(s, col_x + Inches(0.20), y + Inches(0.06), node_w, Inches(0.30), head, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
            add_text(s, col_x + Inches(0.20), y + Inches(0.34), node_w, Inches(0.28), sub, size=11, color=INK, align=PP_ALIGN.CENTER)
        if i < len(nodes) - 1:
            ay = y + node_h + Inches(0.005)
            ax = col_x + col_w / 2 - Inches(0.08)
            add_arrow(s, ax, ay, Inches(0.16), node_gap - Inches(0.01), color=ACCENT, direction="down")

    # U-loop on the RIGHT: from bottom of last node back up to right of "Inference"
    inf_mid_y = content_top + inference_idx * (node_h + node_gap) + node_h / 2
    last_mid_y = content_top + last_idx * (node_h + node_gap) + node_h / 2
    box_right = col_x + Inches(0.20) + node_w
    track_x = box_right + Inches(0.30)
    add_line(s, box_right, last_mid_y, track_x, last_mid_y, color=ACCENT, weight=1.6)
    add_line(s, track_x, last_mid_y, track_x, inf_mid_y, color=ACCENT, weight=1.6)
    add_line(s, track_x, inf_mid_y, box_right, inf_mid_y, color=ACCENT, weight=1.6, head_at_end=True)

    # SupervisorAgent label to the right of the loop track
    sup_mid = (last_mid_y + inf_mid_y) / 2
    add_text(s, track_x + Inches(0.18), sup_mid - Inches(0.40), Inches(2.0), Inches(0.32), "SupervisorAgent", size=14, bold=True, color=NAVY)
    add_text(s, track_x + Inches(0.18), sup_mid - Inches(0.05), Inches(2.0), Inches(0.55), "orchestrates the\nfull loop", size=11, color=GREY, line_spacing=1.15)

    # Human-in-the-loop tag near the human-approval node
    human_y = content_top + 3 * (node_h + node_gap)
    add_text(
        s, col_x - Inches(2.40), human_y + Inches(0.18), Inches(2.20), Inches(0.30),
        "human-in-the-loop", size=11, italic=True, color=AMBER, align=PP_ALIGN.RIGHT,
    )
    # small connector line from the label to the box
    add_line(
        s,
        col_x - Inches(0.20), human_y + node_h / 2,
        col_x + Inches(0.20), human_y + node_h / 2,
        color=AMBER, weight=1.0,
    )

    # Bottom callout: runs locally
    local_top = Inches(6.85)
    add_box(s, Inches(0.5), local_top, Inches(12.33), Inches(0.45), fill=SOFT, border=NAVY, border_w=0.75)
    add_runs(
        s, Inches(0.7), local_top + Inches(0.07), Inches(12.0), Inches(0.32),
        [
            [
                {"text": "Runs locally:  ", "size": 13, "bold": True, "color": NAVY},
                {"text": "MiniLM models on CPU + a local LLM (Ollama) on the same laptop.  No external API required.", "size": 13, "color": INK},
            ]
        ],
    )

    return s


def main():
    path = Path("reference/presentation/intro_slides.pptx")
    if not path.exists():
        raise SystemExit(f"Cannot find {path}. Run scripts/build_intro_slides.py first.")

    prs = Presentation(path)
    n_slides = len(prs.slides)
    if n_slides < 3:
        raise SystemExit(f"Expected at least 3 slides, found {n_slides}.")

    # Drop the existing "agentic workflow" slide (assumed to be slide 3 / index 2)
    delete_slide(prs, 2)

    # Append the two new slides
    s_hitl = build_hitl_slide(prs)
    s_agent = build_agentic_slide(prs)

    # Speaker notes (best-effort -- some master/notes layouts may not expose a
    # notes placeholder; in that case we silently skip).
    def _set_notes(slide, text):
        try:
            tf = slide.notes_slide.notes_text_frame
            if tf is None:
                return
            tf.clear()
            tf.text = text
        except Exception:
            pass

    _set_notes(
        s_hitl,
        "Slide 3 - Human-in-the-loop today (~1 min).\n\n"
        "We start with synthetic data, train both models, run inference, and route pass + borderline cases to human review. "
        "Approved cases feed data augmentation, then retraining, and the loop repeats.\n"
        "Iteration objective is the diagnostic the loop is steering toward: wrong predictions become correct, and borderline "
        "scores between 0.3 and 0.7 move closer to confident 0 or 1.",
    )
    _set_notes(
        s_agent,
        "Slide 4 - Agentic workflow, runs locally (~1.5 min).\n\n"
        "Each manual step from the previous slide becomes a named agent: TrainerAgent, InferenceAgent, ReviewAgent + Human, "
        "InvestigatorAgent, TrainerAgent. SupervisorAgent orchestrates the loop. The human stays at the approval gate.\n"
        "Emphasis: the entire stack runs locally - MiniLM models on CPU and a local LLM via Ollama. No external API required.",
    )

    prs.save(path)
    print(f"Updated {path}: now has {len(prs.slides)} slides.")


if __name__ == "__main__":
    main()
